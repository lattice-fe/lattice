from pathlib import Path

from app.parsers.base import ParsedDocument, ParsedSection, Parser

_EXTENSIONS = {".pptx"}
_MIME_TYPES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}


class PptxParser(Parser):
    """Slide-aware extraction: title, body text, and speaker notes per slide.
    Chart/image vision-description pass is deferred (would need a headless
    LibreOffice render-to-image dependency) — see plan's close-call #4."""

    name = "pptx"
    version = "1"

    @staticmethod
    def supports(mime_type: str | None, extension: str) -> bool:
        return extension.lower() in _EXTENSIONS or mime_type in _MIME_TYPES

    def parse(self, file_path: Path) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        deck_title = _first_title(prs) or file_path.stem

        sections: list[ParsedSection] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_title = _slide_title(slide) or f"Slide {i}"
            body_parts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            text_parts = [p for p in body_parts if p.strip()]
            if notes:
                text_parts.append(f"[speaker notes] {notes}")
            text = "\n\n".join(text_parts).strip()
            if not text:
                continue

            sections.append(
                ParsedSection(
                    heading_path=[deck_title, slide_title],
                    text=text,
                    slide_number=i,
                    metadata={"has_notes": bool(notes)},
                )
            )

        return ParsedDocument(kind="text", sections=sections, raw_metadata={"deck_title": deck_title})


def _first_title(prs) -> str | None:
    if not prs.slides:
        return None
    return _slide_title(prs.slides[0])


def _slide_title(slide) -> str | None:
    if slide.shapes.title is not None and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    return None
