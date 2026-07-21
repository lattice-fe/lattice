"""Builds small sample files on the fly for parser tests rather than
committing binary fixtures — keeps fixtures easy to read and version."""

import csv
from pathlib import Path


def make_pptx(path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Q3 Planning"
    slide1.placeholders[1].text = "Sprint kickoff deck"

    bullet_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "Goals"
    slide2.placeholders[1].text_frame.text = "Ship the ingestion backbone\nWire up permission scoping"
    notes = slide2.notes_slide
    notes.notes_text_frame.text = "Remember to mention the pointer index rollup."

    prs.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["team", "quarter", "spend"])
    ws.append(["platform", "Q3", 42000])
    ws.append(["growth", "Q3", 18000])

    ws2 = wb.create_sheet("Headcount")
    ws2.append(["team", "count"])
    ws2.append(["platform", 6])

    wb.save(str(path))
    return path


def make_csv(path: Path) -> Path:
    with path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["date", "attendee", "topic"])
        writer.writerow(["2026-07-01", "alice", "roadmap"])
        writer.writerow(["2026-07-02", "bob", "incident review"])
    return path


def make_md(path: Path) -> Path:
    path.write_text(
        "# Sprint Notes\n\n"
        "Intro paragraph about the sprint.\n\n"
        "## Decisions\n\n"
        "We decided to ship the backbone first.\n\n"
        "## Risks\n\n"
        "Embedding dimension migrations are manual.\n",
        encoding="utf-8",
    )
    return path


def make_txt(path: Path) -> Path:
    path.write_text(
        "This is a plain text meeting note with no headings.\n\n"
        "It has a second paragraph so paragraph splitting has something to do.\n",
        encoding="utf-8",
    )
    return path


def make_docx(path: Path) -> Path:
    import docx

    doc = docx.Document()
    doc.add_heading("Design Review", level=1)
    doc.add_paragraph("This document captures the design review outcome.")
    doc.add_heading("Open Questions", level=2)
    doc.add_paragraph("Should the pointer index be per-scope or per-folder?")
    doc.save(str(path))
    return path
