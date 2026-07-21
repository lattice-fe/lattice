"""End-to-end smoke test for the ingestion + indexing backbone, per the plan's
verification scenario. Requires a real, disposable Postgres+pgvector instance
(see tests/integration/conftest.py) — skips automatically if none is reachable.

Uses a .docx fixture in place of .pdf for the PdfDocxParser coverage: both
extensions share the same parser class (unstructured.partition.auto), and
generating a well-formed PDF without an extra dependency (reportlab/fpdf,
neither of which this project otherwise needs) isn't reliable, while
python-docx already ships as an unstructured[docx] dependency.
"""

import shutil
import uuid

import pytest
from sqlalchemy import select

from app.constants import EMBEDDING_DIM
from app.jobs.worker import process_one
from app.models.chunks import Chunk, ChunkEmbedding
from app.models.documents import Document, DocumentVersion
from app.models.principals import Principal
from app.models.tabular import TabularRow
from app.permissions.checks import PermissionDenied
from app.repositories import chunk_repo, document_repo, job_repo, scope_repo, summary_repo
from tests.fixtures.generators import make_docx, make_pptx, make_xlsx


class FakeChatProvider:
    def complete(self, messages, *, system=None, max_tokens=1024):
        content = messages[0]["content"]
        return f"[fake summary, {len(content)} chars of prompt]"


class FakeEmbeddingProvider:
    dimension = EMBEDDING_DIM
    model_name = "fake/test-embedding"

    def embed(self, texts):
        return [[0.001 * (i % 7) for _ in range(self.dimension)] for i, _ in enumerate(texts)]


def drain_queue(db_engine, db_session, max_jobs: int = 200) -> int:
    """Repeatedly claims and runs one job at a time until the queue is empty.
    Safe because each stage enqueues its own follow-up job before returning,
    so this naturally walks the whole scan_source -> parse -> chunk -> embed
    -> summarize chain (and any consolidate jobs enqueued by the test).

    Jobs run in their own Sessions (mirroring the real worker), so db_session
    is expired afterward — with expire_on_commit=False, its identity map
    would otherwise keep serving pre-pipeline attribute values for any row
    the test already loaded, masking exactly the state changes being asserted."""
    processed = 0
    for _ in range(max_jobs):
        worked = process_one(f"test-worker-{uuid.uuid4().hex[:6]}")
        if not worked:
            break
        processed += 1
    db_session.expire_all()
    return processed


@pytest.fixture()
def principals(db_session):
    owner = Principal(type="user", external_id="owner@example.com", display_name="Owner")
    p = Principal(type="user", external_id="p@example.com", display_name="P")
    q = Principal(type="user", external_id="q@example.com", display_name="Q")
    db_session.add_all([owner, p, q])
    db_session.commit()
    return {"owner": owner, "p": p, "q": q}


@pytest.fixture()
def two_scopes(db_session, principals, monkeypatch):
    scope_a = scope_repo.create_scope(db_session, principals["owner"].id, "Team A")
    scope_b = scope_repo.create_scope(db_session, principals["owner"].id, "Team B")
    scope_repo.grant_access(db_session, principals["owner"].id, scope_a.id, principals["p"].id, "write")
    scope_repo.grant_access(db_session, principals["owner"].id, scope_b.id, principals["q"].id, "write")
    db_session.commit()

    monkeypatch.setattr("app.pipeline.stages.summarize.get_understanding_chat_provider", lambda: FakeChatProvider())
    monkeypatch.setattr("app.pipeline.stages.consolidate.get_understanding_chat_provider", lambda: FakeChatProvider())
    monkeypatch.setattr("app.pipeline.stages.embed.get_embedding_provider", lambda: FakeEmbeddingProvider())

    return scope_a, scope_b


def test_ingestion_pipeline_end_to_end(db_engine, db_session, principals, two_scopes, tmp_path):
    scope_a, scope_b = two_scopes
    p, q = principals["p"], principals["q"]

    # --- step 2: filesystem_watch source for scope A over a fixture dir ---
    watch_dir = tmp_path / "team_a_docs"
    watch_dir.mkdir()
    make_pptx(watch_dir / "deck.pptx")
    make_xlsx(watch_dir / "budget.xlsx")
    make_docx(watch_dir / "review.docx")

    from app.models.sources import Source

    source_a = Source(scope_id=scope_a.id, type="filesystem_watch", config={"path": str(watch_dir)}, status="active")
    db_session.add(source_a)
    db_session.commit()
    job_repo.enqueue(db_session, "scan_source", {"source_id": str(source_a.id)})
    db_session.commit()

    # scope B gets similar files via manual upload path (as principal Q, who has no access to A)
    source_b = document_repo.get_or_create_manual_upload_source(db_session, scope_b.id)
    upload_dir = tmp_path / "uploads_b"
    upload_dir.mkdir()
    docx_b = make_docx(upload_dir / "review.docx")
    from app.pipeline.hashing import sha256_file

    document_repo.create_or_update_document(
        db_session,
        principal_id=q.id,
        scope_id=scope_b.id,
        source_id=source_b.id,
        external_ref="review.docx",
        mime_type=None,
        file_extension=".docx",
        content_hash=sha256_file(docx_b),
        size_bytes=docx_b.stat().st_size,
        storage_uri=str(docx_b),
    )
    db_session.commit()

    # --- step 3: drain the queue ---
    processed = drain_queue(db_engine, db_session)
    assert processed > 0

    # --- step 4: assert scope A landed correctly ---
    docs_a = document_repo.list_documents(db_session, p.id, scope_a.id)
    assert len(docs_a) == 3
    assert all(d.status == "ready" for d in docs_a)

    by_ref = {d.external_ref: d for d in docs_a}
    for ref in ("deck.pptx", "review.docx"):
        chunks = chunk_repo.list_chunks(db_session, p.id, scope_a.id, by_ref[ref].id)
        assert len(chunks) > 0
        embeddings = db_session.execute(
            select(ChunkEmbedding).where(ChunkEmbedding.chunk_id.in_([c.id for c in chunks]))
        ).scalars().all()
        assert len(embeddings) == len(chunks)
        assert len(embeddings[0].embedding) == FakeEmbeddingProvider.dimension

    xlsx_doc = by_ref["budget.xlsx"]
    assert chunk_repo.list_chunks(db_session, p.id, scope_a.id, xlsx_doc.id) == []
    tabular_rows = db_session.execute(
        select(TabularRow).where(TabularRow.document_version_id == xlsx_doc.latest_version_id)
    ).scalars().all()
    assert len(tabular_rows) > 0

    for doc in docs_a:
        version = db_session.get(DocumentVersion, doc.latest_version_id)
        summaries = [
            s for s in summary_repo.list_summaries_for_scope(db_session, scope_a.id) if s.document_id == doc.id
        ]
        assert len(summaries) == 1
        assert summaries[0].content_hash == version.content_hash

    # --- force consolidation and check pointer index ---
    job_repo.enqueue(db_session, "consolidate", {"scope_id": str(scope_a.id)})
    db_session.commit()
    drain_queue(db_engine, db_session)
    pointer_a = summary_repo.get_pointer_index(db_session, scope_a.id)
    assert pointer_a is not None
    assert pointer_a.doc_count == 3
    assert pointer_a.rollup_text

    # --- step 5: isolation ---
    docs_a_again = document_repo.list_documents(db_session, p.id, scope_a.id)
    assert {d.external_ref for d in docs_a_again} == {"deck.pptx", "budget.xlsx", "review.docx"}

    with pytest.raises(PermissionDenied):
        document_repo.list_documents(db_session, p.id, scope_b.id)

    # --- step 6: mutate the pptx, re-scan ---
    old_chunks = chunk_repo.list_chunks(db_session, p.id, scope_a.id, by_ref["deck.pptx"].id)
    old_chunk_ids = {c.id for c in old_chunks}

    import time

    from pptx import Presentation

    prs = Presentation(str(watch_dir / "deck.pptx"))
    prs.slides[0].shapes.title.text = "Q3 Planning — REVISED"
    time.sleep(0.01)
    prs.save(str(watch_dir / "deck.pptx"))

    job_repo.enqueue(db_session, "scan_source", {"source_id": str(source_a.id)})
    db_session.commit()
    drain_queue(db_engine, db_session)

    refreshed_doc = document_repo.get_document(db_session, p.id, scope_a.id, by_ref["deck.pptx"].id)
    new_chunks = chunk_repo.list_chunks(db_session, p.id, scope_a.id, refreshed_doc.id)
    new_chunk_ids = {c.id for c in new_chunks}
    assert new_chunk_ids.isdisjoint(old_chunk_ids), "old chunks should be replaced, not accumulated"
    assert any("REVISED" in c.content for c in new_chunks)

    job_repo.enqueue(db_session, "consolidate", {"scope_id": str(scope_a.id)})
    db_session.commit()
    drain_queue(db_engine, db_session)
    pointer_after_mutation = summary_repo.get_pointer_index(db_session, scope_a.id)
    assert pointer_after_mutation.doc_count == 3  # still 3 docs, content updated not duplicated

    # --- step 7: delete a file, re-scan ---
    (watch_dir / "budget.xlsx").unlink()
    job_repo.enqueue(db_session, "scan_source", {"source_id": str(source_a.id)})
    db_session.commit()
    drain_queue(db_engine, db_session)

    deleted_doc = db_session.get(Document, xlsx_doc.id)
    assert deleted_doc.deleted_at is not None
    remaining_chunks = db_session.execute(select(Chunk).where(Chunk.document_id == xlsx_doc.id)).scalars().all()
    assert remaining_chunks == []

    job_repo.enqueue(db_session, "consolidate", {"scope_id": str(scope_a.id)})
    db_session.commit()
    drain_queue(db_engine, db_session)
    pointer_final = summary_repo.get_pointer_index(db_session, scope_a.id)
    assert pointer_final.doc_count == 2  # deleted doc excluded

    shutil.rmtree(watch_dir, ignore_errors=True)
